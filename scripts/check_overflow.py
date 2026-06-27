#!/usr/bin/env python3
"""Detect content overflowing the slide canvas in a PPTX file.

Adds gray padding around every slide, renders to images, then inspects
the padding margins pixel-by-pixel to find slides where content bleeds
outside the original canvas boundaries.

Usage:
    python scripts/check_overflow.py output.pptx
    python scripts/check_overflow.py output.pptx --width 1920 --height 1080

Dependencies: pip install python-pptx pdf2image Pillow numpy
              LibreOffice (soffice) + Poppler (pdftoppm)
"""

import argparse
import sys
import tempfile
from os.path import abspath, join
from pathlib import Path
from typing import Sequence, cast

import numpy as np

# Import our render_slides module
SCRIPT_DIR = Path(__file__).resolve().parent
if str(SCRIPT_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPT_DIR))

import render_slides
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Emu

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")

PAD_PX: int = 100
PAD_RGB = (200, 200, 200)
EMU_PER_INCH: int = 914_400


def _px_to_emu(px: int, dpi: int) -> Emu:
    return Emu(int(px * EMU_PER_INCH // dpi))


def _enlarge_deck(src: str, dst: str, pad_emu: Emu) -> tuple[int, int]:
    """Copy the deck with gray padding added around every slide."""
    prs = Presentation(src)
    w0 = cast(Emu, prs.slide_width)
    h0 = cast(Emu, prs.slide_height)
    w1 = Emu(w0 + 2 * pad_emu)
    h1 = Emu(h0 + 2 * pad_emu)
    prs.slide_width = w1
    prs.slide_height = h1

    for slide in prs.slides:
        # Shift all shapes so the original canvas is centered
        for shp in list(slide.shapes):
            shp.left = Emu(int(shp.left) + pad_emu)
            shp.top = Emu(int(shp.top) + pad_emu)

        # Add gray padding rectangles behind everything
        pads = (
            (Emu(0), Emu(0), pad_emu, h1),  # left
            (Emu(int(w1) - int(pad_emu)), Emu(0), pad_emu, h1),  # right
            (Emu(0), Emu(0), w1, pad_emu),  # top
            (Emu(0), Emu(int(h1) - int(pad_emu)), w1, pad_emu),  # bottom
        )
        sp_tree = slide.shapes._spTree
        for left, top, width, height in pads:
            pad_shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height
            )
            pad_shape.fill.solid()
            pad_shape.fill.fore_color.rgb = RGBColor(*PAD_RGB)
            pad_shape.line.fill.background()
            sp_tree.remove(pad_shape._element)
            sp_tree.insert(2, pad_shape._element)

    prs.save(dst)
    return int(w1), int(h1)


def _inspect_margins(
    paths: Sequence[str], pad_ratio_w: float, pad_ratio_h: float, dpi: int
) -> list[int]:
    """Return 1-based indices of slides with content outside the padding."""
    tol = max(1, round((300 - dpi) / 25)) if dpi < 300 else 0
    tol = min(tol, 10)
    pad_colour = np.array(PAD_RGB, dtype=np.uint8)
    failures: list[int] = []

    for idx, img_path in enumerate(paths, start=1):
        with Image.open(img_path) as img:
            arr = np.asarray(img.convert("RGB"))

        h, w, _ = arr.shape
        pad_x = int(w * pad_ratio_w) - 1
        pad_y = int(h * pad_ratio_h) - 1

        margins = [
            arr[:, :pad_x, :],  # left
            arr[:, w - pad_x :, :],  # right
            arr[:pad_y, :, :],  # top
            arr[h - pad_y :, :, :],  # bottom
        ]

        max_mismatch = 0.01 if dpi >= 300 else (0.02 if dpi >= 200 else 0.03)

        for margin in margins:
            diff = np.abs(margin.astype(np.int16) - pad_colour)
            matches = np.all(diff <= tol, axis=-1)
            mismatch = 1.0 - (np.count_nonzero(matches) / matches.size)
            if mismatch > max_mismatch:
                failures.append(idx)
                break

    return failures


def check_overflow(input_path: str, width: int = 1600, height: int = 900) -> list[int]:
    """Check a PPTX for content overflow. Returns list of failing slide numbers."""
    input_path = abspath(input_path)
    dpi = render_slides.calc_dpi(input_path, width, height)

    tmpdir = tempfile.mkdtemp(prefix="overflow_check_")
    enlarged_pptx = join(tmpdir, "enlarged.pptx")
    pad_emu = _px_to_emu(PAD_PX, dpi)
    w1, h1 = _enlarge_deck(input_path, enlarged_pptx, pad_emu)
    pad_ratio_w = pad_emu / w1
    pad_ratio_h = pad_emu / h1

    img_dir = join(tmpdir, "imgs")
    img_paths = render_slides.rasterize(enlarged_pptx, img_dir, dpi)
    return _inspect_margins(img_paths, pad_ratio_w, pad_ratio_h, dpi)


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Check PPTX for content overflowing the slide canvas."
    )
    parser.add_argument("input_path", help="Path to PPTX file")
    parser.add_argument(
        "--width", type=int, default=1600, help="Target render width in pixels (default: 1600)"
    )
    parser.add_argument(
        "--height", type=int, default=900, help="Target render height in pixels (default: 900)"
    )
    args = parser.parse_args()

    failing = check_overflow(args.input_path, args.width, args.height)
    if failing:
        print(
            f"Overflow detected on {len(failing)} slide(s): " + ", ".join(str(s) for s in failing)
        )
        sys.exit(1)
    else:
        print("No overflow detected")


if __name__ == "__main__":
    main()
