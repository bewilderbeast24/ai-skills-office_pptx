#!/usr/bin/env python3
"""Render PPTX slides to individual PNG images with auto-calculated DPI.

Smarter than raw soffice+pdftoppm: reads the actual slide dimensions from
the OOXML XML to compute DPI that hits a target pixel resolution, and has
a PPTX→ODP→PDF fallback for problematic decks.

Usage:
    python scripts/render_slides.py output.pptx
    python scripts/render_slides.py output.pptx --output_dir slides/
    python scripts/render_slides.py output.pptx --width 1920 --height 1080

Dependencies: pip install pdf2image Pillow
              LibreOffice (soffice) + Poppler (pdftoppm)
"""

import argparse
import shutil
import subprocess
import tempfile
import xml.etree.ElementTree as ET
from os import makedirs, replace
from os.path import abspath, basename, exists, expanduser, join, splitext
from typing import Sequence, cast
from zipfile import ZipFile

from pdf2image import convert_from_path
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

# Import our sandboxed soffice helper
try:
    from office.soffice import get_soffice_env
except ImportError:
    import sys

    sys.path.insert(0, str(__import__("pathlib").Path(__file__).resolve().parent))
    from office.soffice import get_soffice_env

EMU_PER_INCH: int = 914_400


def _run_soffice(args: list[str], outdir: str, user_profile: str) -> None:
    """Run soffice with our sandbox-aware env and a unique user profile."""
    env = get_soffice_env()
    profile_uri = __import__("pathlib").Path(user_profile).resolve().as_uri()
    cmd = [
        "soffice",
        f"-env:UserInstallation={profile_uri}",
        "--invisible",
        "--headless",
        "--norestore",
    ] + args
    subprocess.run(cmd, check=False, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, env=env)


def calc_dpi(input_path: str, max_w_px: int, max_h_px: int) -> int:
    """Calculate DPI from OOXML slide size to hit target pixel dimensions."""
    with ZipFile(input_path, "r") as zf:
        xml = zf.read("ppt/presentation.xml")
    root = ET.fromstring(xml)
    ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    sld_sz = root.find("p:sldSz", ns)
    if sld_sz is None:
        return 150  # fallback
    cx = int(sld_sz.get("cx") or 0)
    cy = int(sld_sz.get("cy") or 0)
    if cx <= 0 or cy <= 0:
        return 150
    width_in = cx / EMU_PER_INCH
    height_in = cy / EMU_PER_INCH
    return round(min(max_w_px / width_in, max_h_px / height_in))


def _convert_to_pdf(pptx_path: str, user_profile: str, out_dir: str, stem: str) -> str:
    """Convert PPTX to PDF, with ODP fallback for problematic files."""
    # Try direct PPTX → PDF
    _run_soffice(["--convert-to", "pdf", "--outdir", out_dir, pptx_path], out_dir, user_profile)
    pdf_path = join(out_dir, f"{stem}.pdf")
    if exists(pdf_path):
        return pdf_path

    # Fallback: PPTX → ODP → PDF
    _run_soffice(["--convert-to", "odp", "--outdir", out_dir, pptx_path], out_dir, user_profile)
    odp_path = join(out_dir, f"{stem}.odp")
    if exists(odp_path):
        _run_soffice(["--convert-to", "pdf", "--outdir", out_dir, odp_path], out_dir, user_profile)
        if exists(pdf_path):
            return pdf_path
    return ""


def rasterize(input_path: str, out_dir: str, dpi: int) -> Sequence[str]:
    """Rasterize PPTX/PDF to numbered PNG files and return their paths."""
    makedirs(out_dir, exist_ok=True)
    input_path = abspath(input_path)
    stem = splitext(basename(input_path))[0]

    if not input_path.lower().endswith(".pdf") and not _find_soffice():
        return _rasterize_with_python_pptx(input_path, out_dir, dpi)

    with tempfile.TemporaryDirectory(prefix="soffice_profile_") as user_profile:
        with tempfile.TemporaryDirectory(prefix="soffice_convert_") as convert_dir:
            is_pdf = input_path.lower().endswith(".pdf")
            pdf_path = (
                input_path
                if is_pdf
                else _convert_to_pdf(input_path, user_profile, convert_dir, stem)
            )

            if not pdf_path or not exists(pdf_path):
                if not is_pdf:
                    return _rasterize_with_python_pptx(input_path, out_dir, dpi)
                raise RuntimeError("Failed to produce PDF for rasterization.")

            try:
                paths_raw = cast(
                    list[str],
                    convert_from_path(
                        pdf_path,
                        dpi=dpi,
                        fmt="png",
                        thread_count=8,
                        output_folder=out_dir,
                        paths_only=True,
                        output_file="slide",
                    ),
                )
            except Exception:
                paths_raw = list(_rasterize_pdf_with_pypdfium(pdf_path, out_dir, dpi))

    # Rename to clean slide-1.png, slide-2.png, ... format
    slides = []
    for src_path in paths_raw:
        base = splitext(basename(src_path))[0]
        slide_num_str = base.split("-")[-1]
        slide_num = int(slide_num_str)
        dst_path = join(out_dir, f"slide-{slide_num}.png")
        replace(src_path, dst_path)
        slides.append((slide_num, dst_path))
    slides.sort(key=lambda t: t[0])
    return [path for _, path in slides]


def _find_soffice() -> str | None:
    found = shutil.which("soffice")
    if found:
        return found
    for candidate in (
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ):
        if exists(candidate):
            return candidate
    return None


def _rasterize_pdf_with_pypdfium(pdf_path: str, out_dir: str, dpi: int) -> Sequence[str]:
    import pypdfium2 as pdfium

    pdf = pdfium.PdfDocument(pdf_path)
    scale = dpi / 72
    paths: list[str] = []
    try:
        for idx in range(len(pdf)):
            page = pdf[idx]
            try:
                bitmap = page.render(scale=scale)
                try:
                    image = bitmap.to_pil()
                    path = join(out_dir, f"slide-pdfium-{idx + 1}.png")
                    image.save(path)
                    paths.append(path)
                finally:
                    if hasattr(bitmap, "close"):
                        bitmap.close()
            finally:
                if hasattr(page, "close"):
                    page.close()
    finally:
        if hasattr(pdf, "close"):
            pdf.close()
    return paths


def _rgb_to_hex(rgb, default: str) -> str:
    try:
        return f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
    except Exception:
        return default


def _shape_fill(shape, default: str = "#F2F2F2") -> str:
    try:
        fill = shape.fill
        if fill.type is None:
            return default
        return _rgb_to_hex(fill.fore_color.rgb, default)
    except Exception:
        return default


def _text_color(run, default: str = "#111111") -> str:
    try:
        return _rgb_to_hex(run.font.color.rgb, default)
    except Exception:
        return default


def _font(size: int):
    try:
        return ImageFont.truetype("arial.ttf", max(8, size))
    except Exception:
        return ImageFont.load_default()


def _rasterize_with_python_pptx(input_path: str, out_dir: str, dpi: int) -> Sequence[str]:
    """Fallback visual overview renderer used when LibreOffice is unavailable.

    This is not a pixel-perfect PowerPoint renderer. It draws common shapes,
    text, and chart placeholders well enough for thumbnail grids and overflow
    checks to keep working in lightweight environments.
    """
    prs = Presentation(input_path)
    width_px = max(1, round(int(prs.slide_width) / EMU_PER_INCH * dpi))
    height_px = max(1, round(int(prs.slide_height) / EMU_PER_INCH * dpi))
    sx = width_px / int(prs.slide_width)
    sy = height_px / int(prs.slide_height)

    paths: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        bg = "#FFFFFF"
        try:
            bg = _rgb_to_hex(slide.background.fill.fore_color.rgb, bg)
        except Exception:
            pass
        img = Image.new("RGB", (width_px, height_px), bg)
        draw = ImageDraw.Draw(img)

        for shape in slide.shapes:
            x = round(int(shape.left) * sx)
            y = round(int(shape.top) * sy)
            w = round(int(shape.width) * sx)
            h = round(int(shape.height) * sy)
            if w <= 0 or h <= 0:
                continue

            fill = _shape_fill(shape)
            outline = "#777777"
            try:
                outline = _rgb_to_hex(shape.line.color.rgb, outline)
            except Exception:
                pass

            if getattr(shape, "has_chart", False):
                draw.rectangle([x, y, x + w, y + h], fill="#F8F8F8", outline=outline, width=2)
                draw.text((x + 12, y + 12), "Chart", fill="#333333", font=_font(18))
            elif not getattr(shape, "has_text_frame", False) or not shape.text.strip():
                draw.rectangle([x, y, x + w, y + h], fill=fill, outline=outline)

            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                tx = x + 8
                ty = y + 6
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs) or paragraph.text
                    if not line:
                        ty += 14
                        continue
                    first_run = paragraph.runs[0] if paragraph.runs else None
                    size = 16
                    color = "#111111"
                    if first_run is not None:
                        try:
                            if first_run.font.size:
                                size = max(8, round(first_run.font.size.pt * dpi / 96))
                        except Exception:
                            pass
                        color = _text_color(first_run, color)
                    font = _font(size)
                    draw.multiline_text((tx, ty), line, fill=color, font=font, spacing=4)
                    bbox = draw.multiline_textbbox((tx, ty), line, font=font, spacing=4)
                    ty = bbox[3] + 6

        path = join(out_dir, f"slide-{idx}.png")
        img.save(path)
        paths.append(path)

    return paths


def main() -> None:
    parser = argparse.ArgumentParser(description="Render PPTX/PDF slides to PNG images.")
    parser.add_argument("input_path", help="Path to PPTX or PDF file")
    parser.add_argument(
        "--output_dir",
        default=None,
        help="Output directory (default: folder named after input file)",
    )
    parser.add_argument(
        "--width", type=int, default=1600, help="Target max width in pixels (default: 1600)"
    )
    parser.add_argument(
        "--height", type=int, default=900, help="Target max height in pixels (default: 1600)"
    )
    args = parser.parse_args()

    input_path = abspath(expanduser(args.input_path))
    out_dir = abspath(expanduser(args.output_dir)) if args.output_dir else splitext(input_path)[0]
    dpi = calc_dpi(input_path, args.width, args.height)
    paths = rasterize(input_path, out_dir, dpi)
    print(f"Rendered {len(paths)} slides to {out_dir}")


if __name__ == "__main__":
    main()
